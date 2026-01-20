"""
PowerPoint Generator Tool for E-Powertrain Benchmarking System

This tool transforms scraped vehicle data into IAA template format
and generates PowerPoint presentations.

Workflow:
1. Receive ScrapingResult or VehicleSpecifications from state
2. Transform to OEMProfile (IAA template schema)
3. Populate IAA template
4. Return PresentationResult
"""

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from typing import Dict, List, Optional, Any
from datetime import datetime
from pathlib import Path
import json
import re
import time


# =====================================================================
# SHAPE ID MAPPING - Discovered from IAA template analysis
# =====================================================================

SHAPE_IDS = {
    # Main data entry points
    "company_name": 2,           # Title placeholder - Company name
    "company_info_table": 7,     # Table: Company Information (6x2)
    "products_table": 8,         # Table: Product Specifications (14x3)
    "expected_highlights": 25,   # TextBox: Expected Highlights bullets
    "assessment": 27,            # TextBox: Assessment bullets
    "technologies": 26,          # TextBox: Technologies bullets
    
    # Headers (static)
    "oem_label": 23,
    "cooperations_header": 17,
    
    # Content areas
    "pictures_area": 20,
    "cooperations_area": 21,
    "logo_area": 22,
    
    # Footer
    "footer": 4,
    "slide_number": 5,
}

# Product table row indices
PRODUCT_TABLE_ROWS = {
    "name": 0,
    "wheel_formula": 1,
    "wheelbase": 2,
    "gvw_gcw": 3,
    "range": 4,
    "battery": 5,
    "fuel_cell": 6,
    "h2_tank": 7,
    "charging": 8,
    "performance": 9,
    "powertrain": 10,
    "sop": 11,
    "markets": 12,
    "application": 13,
}


# =====================================================================
# DATA TRANSFORMATION: VehicleSpecifications -> IAA_ProductSpec
# =====================================================================

def transform_vehicle_to_iaa_product(vehicle: Dict) -> Dict:
    """
    Transform a VehicleSpecifications dict to IAA_ProductSpec format.
    Handles unit conversions and formatting.
    """
    
    def format_value(value, unit: str = "") -> str:
        """Format a value with unit, handling None/missing"""
        if value is None:
            return "N/A"
        if isinstance(value, (int, float)):
            if value == 0:
                return "N/A"
            return f"{value:,.0f} {unit}".strip() if value == int(value) else f"{value:,.1f} {unit}".strip()
        return str(value) if value else "N/A"
    
    # Get additional specs if available
    additional = vehicle.get('additional_specs', {}) or {}
    
    return {
        "name": vehicle.get('vehicle_name', 'Unknown Vehicle'),
        "wheel_formula": additional.get('wheel_formula', 'N/A'),
        "wheelbase": format_value(additional.get('wheelbase_mm'), 'mm'),
        "gvw_gcw": format_value(vehicle.get('gvw_kg'), 'kg GVW'),
        "range": format_value(vehicle.get('range_km'), 'km'),
        "battery": format_value(vehicle.get('battery_capacity_kwh'), 'kWh'),
        "fuel_cell": format_value(additional.get('fuel_cell_kw'), 'kW'),
        "h2_tank": format_value(additional.get('h2_tank_kg'), 'kg'),
        "charging": format_value(vehicle.get('dc_charging_kw'), 'kW DC'),
        "performance": format_value(vehicle.get('motor_power_kw'), 'kW'),
        "powertrain": additional.get('powertrain_description', 
                      vehicle.get('powertrain_type', 'Electric')),
        "sop": additional.get('start_of_production', 'TBD'),
        "markets": additional.get('markets', 'N/A'),
        "application": vehicle.get('category', 'Commercial Vehicle'),
    }


def transform_scraping_result_to_oem_profile(scraping_result: Dict) -> Dict:
    """
    Transform a ScrapingResult to OEMProfile format for IAA template.
    
    This is the main transformation function that bridges the gap between
    the scraper output and the presentation template.
    """
    oem_name = scraping_result.get('oem_name', 'Unknown OEM')
    oem_url = scraping_result.get('oem_url', '')
    vehicles = scraping_result.get('vehicles', [])
    
    # Extract domain for website
    domain = oem_url.split('/')[2] if '://' in oem_url else oem_url
    
    # Transform vehicles to IAA product format
    products = [transform_vehicle_to_iaa_product(v) for v in vehicles[:2]]
    
    # Generate highlights from vehicle data
    highlights = generate_highlights(vehicles, oem_name)
    
    # Generate assessment based on data quality
    assessment = generate_assessment(scraping_result)
    
    # Extract technologies from vehicles
    technologies = extract_technologies(vehicles)
    
    return {
        "company_name": oem_name,
        "company_info": {
            "country": extract_country_from_url(oem_url),
            "address": "",  # Would need additional scraping
            "website": domain,
            "booth": "",    # Event-specific, may be input
            "category": determine_oem_category(vehicles),
        },
        "expected_highlights": highlights,
        "assessment": assessment,
        "technologies": technologies,
        "products": products,
        "cooperations": [],  # Would need additional scraping
        "source_url": oem_url,
        "extraction_timestamp": scraping_result.get('extraction_timestamp', 
                                                     datetime.now().isoformat()),
        "data_quality_score": scraping_result.get('source_compliance_score', 0.0),
    }


def generate_highlights(vehicles: List[Dict], oem_name: str) -> List[str]:
    """Generate expected highlights from vehicle data"""
    highlights = []
    
    if vehicles:
        # Count by powertrain type
        bev_count = sum(1 for v in vehicles if 'bev' in str(v.get('powertrain_type', '')).lower() 
                        or 'battery' in str(v.get('powertrain_type', '')).lower()
                        or v.get('battery_capacity_kwh'))
        fcev_count = sum(1 for v in vehicles if 'fuel cell' in str(v.get('powertrain_type', '')).lower()
                         or v.get('additional_specs', {}).get('fuel_cell_kw'))
        
        if bev_count > 0:
            highlights.append(f"{bev_count} Battery Electric Vehicle(s) in portfolio")
        if fcev_count > 0:
            highlights.append(f"{fcev_count} Fuel Cell Vehicle(s) available")
        
        # Range highlights
        max_range = max((v.get('range_km', 0) or 0 for v in vehicles), default=0)
        if max_range > 0:
            highlights.append(f"Up to {max_range:.0f} km range capability")
        
        # Battery highlights
        max_battery = max((v.get('battery_capacity_kwh', 0) or 0 for v in vehicles), default=0)
        if max_battery > 0:
            highlights.append(f"Battery capacity up to {max_battery:.0f} kWh")
    
    # Default if no highlights generated
    if not highlights:
        highlights = [f"{oem_name} e-mobility portfolio", "Technical specifications available"]
    
    return highlights[:4]  # Limit to 4 highlights


def generate_assessment(scraping_result: Dict) -> List[str]:
    """Generate assessment based on data quality and coverage"""
    assessment = []
    
    score = scraping_result.get('source_compliance_score', 0)
    vehicles = scraping_result.get('vehicles', [])
    official_citations = len(scraping_result.get('official_citations', []))
    
    if score >= 0.8:
        assessment.append("High data quality from official sources")
    elif score >= 0.5:
        assessment.append("Moderate data quality, some gaps")
    else:
        assessment.append("Limited data available, verification needed")
    
    if len(vehicles) >= 3:
        assessment.append(f"Comprehensive portfolio with {len(vehicles)} vehicles")
    elif len(vehicles) >= 1:
        assessment.append(f"{len(vehicles)} vehicle(s) documented")
    
    if official_citations >= 2:
        assessment.append("Multiple official source confirmations")
    
    return assessment[:3]


def extract_technologies(vehicles: List[Dict]) -> List[str]:
    """Extract technology keywords from vehicle data"""
    technologies = set()
    
    for v in vehicles:
        if v.get('battery_capacity_kwh'):
            chemistry = v.get('battery_chemistry', '')
            if chemistry:
                technologies.add(f"{chemistry} Battery")
            else:
                technologies.add("Li-ion Battery")
        
        if v.get('dc_charging_kw'):
            dc_power = v.get('dc_charging_kw')
            if dc_power and dc_power >= 150:
                technologies.add(f"High-Power DC Charging ({dc_power:.0f}kW)")
            elif dc_power:
                technologies.add("DC Fast Charging")
        
        additional = v.get('additional_specs', {}) or {}
        if additional.get('fuel_cell_kw'):
            technologies.add("Fuel Cell System")
        
        powertrain = v.get('powertrain_type', '')
        if 'hybrid' in powertrain.lower():
            technologies.add("Hybrid Powertrain")
    
    return list(technologies)[:5]


def extract_country_from_url(url: str) -> str:
    """Guess country from URL TLD"""
    tld_country = {
        '.de': 'Germany', '.eu': 'Europe', '.com': 'USA',
        '.cn': 'China', '.jp': 'Japan', '.kr': 'South Korea',
        '.se': 'Sweden', '.nl': 'Netherlands', '.fr': 'France',
        '.it': 'Italy', '.uk': 'United Kingdom',
    }
    for tld, country in tld_country.items():
        if tld in url:
            return country
    return ""


def determine_oem_category(vehicles: List[Dict]) -> str:
    """Determine OEM category based on vehicle types"""
    categories = [v.get('category', '').lower() for v in vehicles if v.get('category')]
    
    if any('truck' in c for c in categories):
        return "OEM - Commercial Trucks"
    elif any('bus' in c for c in categories):
        return "OEM - Bus & Coach"
    elif any('van' in c for c in categories):
        return "OEM - Light Commercial"
    return "OEM - Commercial Vehicles"


# =====================================================================
# POWERPOINT GENERATION
# =====================================================================

def get_shape_by_id(slide, shape_id: int):
    """Find a shape by its shape_id"""
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    return None


def set_bullet_list(shape, items: List[str], font_size: int = 9):
    """Populate a shape with bullet point list"""
    if not hasattr(shape, "text_frame") or not items:
        return
    
    tf = shape.text_frame
    
    # Clear existing
    for p in tf.paragraphs:
        p.clear()
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.level = 0
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)


def populate_company_info_table(table, company_info: Dict):
    """Populate the Company Information table (6x2)"""
    field_mapping = {
        1: "country",
        2: "address",
        3: "website",
        4: "booth",
        5: "category"
    }
    
    for row_idx, field_name in field_mapping.items():
        value = company_info.get(field_name, "")
        if value:
            cell = table.cell(row_idx, 1)
            cell.text = str(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(9)


def populate_products_table(table, products: List[Dict]):
    """Populate the Products Specifications table (14x3)"""
    for product_idx, product in enumerate(products[:2]):
        col = product_idx + 1
        
        for field_name, row_idx in PRODUCT_TABLE_ROWS.items():
            value = product.get(field_name, "")
            if value:
                cell = table.cell(row_idx, col)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(8)


def generate_presentation(
    oem_profile: Dict,
    template_path: str,
    output_path: str
) -> Dict:
    """
    Generate PowerPoint presentation from OEM profile.
    
    Args:
        oem_profile: OEMProfile dict with all data
        template_path: Path to IAA template .pptx
        output_path: Output path for generated .pptx
    
    Returns:
        PresentationResult dict
    """
    start_time = time.time()
    
    prs = Presentation(template_path)
    slide = prs.slides[0]
    
    # 1. Company Name
    shape = get_shape_by_id(slide, SHAPE_IDS["company_name"])
    if shape:
        shape.text = oem_profile.get("company_name", "Unknown OEM")
    
    # 2. Company Information Table
    if oem_profile.get("company_info"):
        shape = get_shape_by_id(slide, SHAPE_IDS["company_info_table"])
        if shape and shape.has_table:
            populate_company_info_table(shape.table, oem_profile["company_info"])
    
    # 3. Products Table
    if oem_profile.get("products"):
        shape = get_shape_by_id(slide, SHAPE_IDS["products_table"])
        if shape and shape.has_table:
            populate_products_table(shape.table, oem_profile["products"])
    
    # 4. Expected Highlights
    if oem_profile.get("expected_highlights"):
        shape = get_shape_by_id(slide, SHAPE_IDS["expected_highlights"])
        if shape:
            set_bullet_list(shape, oem_profile["expected_highlights"])
    
    # 5. Assessment
    if oem_profile.get("assessment"):
        shape = get_shape_by_id(slide, SHAPE_IDS["assessment"])
        if shape:
            set_bullet_list(shape, oem_profile["assessment"])
    
    # 6. Technologies
    if oem_profile.get("technologies"):
        shape = get_shape_by_id(slide, SHAPE_IDS["technologies"])
        if shape:
            set_bullet_list(shape, oem_profile["technologies"])
    
    # Save
    prs.save(output_path)
    
    duration = time.time() - start_time
    
    return {
        "presentation_path": output_path,
        "slides_created": 1,
        "vehicles_included": len(oem_profile.get("products", [])),
        "oems_compared": 1,
        "includes_charts": False,
        "includes_comparison_table": True,
        "generation_timestamp": datetime.now().isoformat(),
        "generation_duration_seconds": duration,
    }


# =====================================================================
# MAIN TOOL FUNCTION - Called by Presentation Agent
# =====================================================================

def generate_oem_presentation(
    scraping_result: Dict,
    template_path: str = "templates/IAA_Template.pptx",
    output_dir: str = "outputs"
) -> Dict:
    """
    Main entry point for presentation generation.
    
    Called by the Presentation Generator Agent with ScrapingResult.
    
    Args:
        scraping_result: ScrapingResult dict from scraper agent
        template_path: Path to IAA template
        output_dir: Directory for output files
    
    Returns:
        Dict with 'oem_profile' and 'presentation_result'
    """
    # Ensure output directory exists
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    
    # Transform to OEM profile
    oem_profile = transform_scraping_result_to_oem_profile(scraping_result)
    
    # Generate safe filename
    safe_name = re.sub(r'[^\w\-]', '_', oem_profile['company_name'])
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = f"{output_dir}/IAA_{safe_name}_{timestamp}.pptx"
    
    # Generate presentation
    result = generate_presentation(oem_profile, template_path, output_path)
    
    return {
        "oem_profile": oem_profile,
        "presentation_result": result,
    }


# =====================================================================
# BATCH PROCESSING - For multiple OEMs
# =====================================================================

def generate_all_presentations(
    scraping_results: List[Dict],
    template_path: str = "templates/IAA_Template.pptx",
    output_dir: str = "outputs"
) -> Dict:
    """
    Generate presentations for ALL OEMs in batch.

    This is the main entry point for the Presentation Generator Agent.
    Processes all ScrapingResults and returns aggregated results
    matching the BenchmarkingState schema.

    Args:
        scraping_results: List of ScrapingResult dicts from scraper
        template_path: Path to IAA template
        output_dir: Directory for output files

    Returns:
        Dict with:
            - oem_profiles: List[OEMProfile] for state.oem_profiles
            - presentation_result: PresentationResult (aggregated)
            - individual_results: List of per-OEM results
    """
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    start_time = time.time()
    oem_profiles = []
    individual_results = []
    presentation_paths = []
    total_vehicles = 0
    total_slides = 0
    errors = []

    for scraping_result in scraping_results:
        try:
            # Transform to OEM profile
            oem_profile = transform_scraping_result_to_oem_profile(scraping_result)
            oem_profiles.append(oem_profile)

            # Generate safe filename
            safe_name = re.sub(r'[^\w\-]', '_', oem_profile['company_name'])
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = f"{output_dir}/IAA_{safe_name}_{timestamp}.pptx"

            # Generate presentation
            result = generate_presentation(oem_profile, template_path, output_path)

            individual_results.append({
                "oem_name": oem_profile['company_name'],
                "oem_profile": oem_profile,
                "presentation_result": result,
            })

            presentation_paths.append(output_path)
            total_vehicles += result.get("vehicles_included", 0)
            total_slides += result.get("slides_created", 0)

        except Exception as e:
            error_msg = f"Failed to generate presentation for {scraping_result.get('oem_name', 'Unknown')}: {str(e)}"
            errors.append(error_msg)
            print(f"ERROR: {error_msg}")

    total_duration = time.time() - start_time

    # Create aggregated PresentationResult matching state schema
    aggregated_result = {
        "presentation_path": presentation_paths[0] if presentation_paths else "",
        "all_presentation_paths": presentation_paths,  # Extra: all paths
        "slides_created": total_slides,
        "vehicles_included": total_vehicles,
        "oems_compared": len(oem_profiles),
        "includes_charts": False,
        "includes_comparison_table": True,
        "generation_timestamp": datetime.now().isoformat(),
        "generation_duration_seconds": round(total_duration, 2),
        "errors": errors,
    }

    return {
        "oem_profiles": oem_profiles,
        "presentation_result": aggregated_result,
        "individual_results": individual_results,
    }


def generate_comparison_presentation(
    scraping_results: List[Dict],
    template_path: str = "templates/IAA_Template.pptx",
    output_dir: str = "outputs"
) -> Dict:
    """
    Generate a SINGLE comparison presentation with all OEMs.

    Alternative to generate_all_presentations() when you want
    one presentation with multiple OEM slides instead of separate files.

    Args:
        scraping_results: List of ScrapingResult dicts
        template_path: Path to template
        output_dir: Output directory

    Returns:
        Dict with presentation_result and oem_profiles
    """
    Path(output_dir).mkdir(parents=True, exist_ok=True)
    start_time = time.time()

    # Load template
    prs = Presentation(template_path)
    base_slide = prs.slides[0]

    oem_profiles = []
    total_vehicles = 0

    for i, scraping_result in enumerate(scraping_results):
        oem_profile = transform_scraping_result_to_oem_profile(scraping_result)
        oem_profiles.append(oem_profile)

        if i == 0:
            # Use the first slide from template
            slide = base_slide
        else:
            # Duplicate the template slide for additional OEMs
            slide_layout = base_slide.slide_layout
            slide = prs.slides.add_slide(slide_layout)

        # Populate slide with OEM data
        _populate_slide(slide, oem_profile)
        total_vehicles += len(oem_profile.get("products", []))

    # Save combined presentation
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = f"{output_dir}/IAA_Comparison_{timestamp}.pptx"
    prs.save(output_path)

    duration = time.time() - start_time

    return {
        "oem_profiles": oem_profiles,
        "presentation_result": {
            "presentation_path": output_path,
            "slides_created": len(prs.slides),
            "vehicles_included": total_vehicles,
            "oems_compared": len(oem_profiles),
            "includes_charts": False,
            "includes_comparison_table": True,
            "generation_timestamp": datetime.now().isoformat(),
            "generation_duration_seconds": round(duration, 2),
        },
    }


def _populate_slide(slide, oem_profile: Dict):
    """Helper to populate a single slide with OEM data"""

    # 1. Company Name
    shape = get_shape_by_id(slide, SHAPE_IDS["company_name"])
    if shape:
        shape.text = oem_profile.get("company_name", "Unknown OEM")

    # 2. Company Information Table
    if oem_profile.get("company_info"):
        shape = get_shape_by_id(slide, SHAPE_IDS["company_info_table"])
        if shape and shape.has_table:
            populate_company_info_table(shape.table, oem_profile["company_info"])

    # 3. Products Table
    if oem_profile.get("products"):
        shape = get_shape_by_id(slide, SHAPE_IDS["products_table"])
        if shape and shape.has_table:
            populate_products_table(shape.table, oem_profile["products"])

    # 4. Expected Highlights
    if oem_profile.get("expected_highlights"):
        shape = get_shape_by_id(slide, SHAPE_IDS["expected_highlights"])
        if shape:
            set_bullet_list(shape, oem_profile["expected_highlights"])

    # 5. Assessment
    if oem_profile.get("assessment"):
        shape = get_shape_by_id(slide, SHAPE_IDS["assessment"])
        if shape:
            set_bullet_list(shape, oem_profile["assessment"])

    # 6. Technologies
    if oem_profile.get("technologies"):
        shape = get_shape_by_id(slide, SHAPE_IDS["technologies"])
        if shape:
            set_bullet_list(shape, oem_profile["technologies"])


# =====================================================================
# CLI / TEST
# =====================================================================

if __name__ == "__main__":
    # Test with mock scraping result
    mock_scraping_result = {
        "oem_name": "MAN Truck & Bus",
        "oem_url": "https://www.man.eu",
        "vehicles": [
            {
                "vehicle_name": "MAN eTGX",
                "oem_name": "MAN",
                "category": "Heavy-duty Truck",
                "powertrain_type": "BEV",
                "source_url": "https://www.man.eu/etgx",
                "battery_capacity_kwh": 480,
                "motor_power_kw": 480,
                "range_km": 600,
                "dc_charging_kw": 375,
                "gvw_kg": 40000,
                "additional_specs": {
                    "wheel_formula": "6x2",
                    "wheelbase_mm": 4500,
                    "powertrain_description": "Central e-motor",
                    "start_of_production": "2024",
                    "markets": "EU",
                },
                "extraction_timestamp": datetime.now().isoformat(),
            },
            {
                "vehicle_name": "MAN eTGM",
                "oem_name": "MAN",
                "category": "Medium-duty Truck",
                "powertrain_type": "BEV",
                "source_url": "https://www.man.eu/etgm",
                "battery_capacity_kwh": 264,
                "motor_power_kw": 264,
                "range_km": 300,
                "dc_charging_kw": 150,
                "gvw_kg": 26000,
                "additional_specs": {
                    "wheel_formula": "4x2",
                    "wheelbase_mm": 3900,
                    "powertrain_description": "E-axle",
                    "start_of_production": "2023",
                    "markets": "EU, UK",
                },
                "extraction_timestamp": datetime.now().isoformat(),
            }
        ],
        "total_vehicles_found": 2,
        "extraction_timestamp": datetime.now().isoformat(),
        "official_citations": [
            "https://www.man.eu/de/de/lkw/etruck/etgx/uebersicht.html",
            "https://www.man.eu/de/de/lkw/etruck/etgm/uebersicht.html"
        ],
        "third_party_citations": [],
        "source_compliance_score": 0.85,
        "raw_content": "...",
        "tokens_used": 1500,
        "model_used": "sonar-pro",
        "extraction_duration_seconds": 12.5,
        "errors": [],
        "warnings": [],
    }
    
    print("Testing PPT Generator...")
    print("=" * 60)
    
    # Run generation (requires template file)
    try:
        result = generate_oem_presentation(
            scraping_result=mock_scraping_result,
            template_path="IAA_Template.pptx",
            output_dir="."
        )
        
        print(f"OEM Profile: {result['oem_profile']['company_name']}")
        print(f"Products: {len(result['oem_profile']['products'])}")
        print(f"Highlights: {result['oem_profile']['expected_highlights']}")
        print(f"Technologies: {result['oem_profile']['technologies']}")
        print(f"\nPresentation saved to: {result['presentation_result']['presentation_path']}")
        
    except FileNotFoundError:
        print("Template file not found. Place IAA_Template.pptx in current directory.")
        
        # Still show transformation
        oem_profile = transform_scraping_result_to_oem_profile(mock_scraping_result)
        print("\nTransformed OEM Profile:")
        print(json.dumps(oem_profile, indent=2, default=str))