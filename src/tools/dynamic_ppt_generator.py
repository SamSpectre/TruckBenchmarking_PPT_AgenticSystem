"""
Dynamic PowerPoint Generator for Plug-and-Play Template System

Generates presentations using JSON mapping configurations.
Supports any template with a valid mapping config.
Supports multi-slide templates with overflow/pagination for arrays.

This is a standalone module - does NOT modify existing ppt_generator.py

Version History:
- 1.0.0: Initial release with single-slide support
- 2.0.0: Added multi-slide support with full slide duplication
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsmap
from typing import Dict, List, Any, Optional, Union, Tuple
from datetime import datetime
from pathlib import Path
from copy import deepcopy
import json
import re
import time
import math


class DynamicPPTGenerator:
    """
    Generates PowerPoint presentations using JSON mapping configurations.
    Supports multi-slide templates with overflow/pagination for arrays.

    Usage:
        generator = DynamicPPTGenerator(mapping_config)
        result = generator.generate(template_path, data, output_path)
    """

    def __init__(self, mapping_config: Union[Dict, str, Path]):
        """
        Initialize generator with a mapping configuration.

        Args:
            mapping_config: Either a dict, path to JSON file, or Path object
        """
        if isinstance(mapping_config, (str, Path)):
            with open(mapping_config, 'r', encoding='utf-8') as f:
                self.mapping = json.load(f)
        else:
            self.mapping = mapping_config

        self.template_name = self.mapping.get("template_name", "unknown")
        self.supports_multi_slide = self.mapping.get("supports_multi_slide", False)
        self.pagination_rules = self.mapping.get("pagination_rules", {})

    def generate(
        self,
        template_path: str,
        data: Dict[str, Any],
        output_path: str
    ) -> Dict[str, Any]:
        """
        Generate a presentation from template using data.
        Supports multi-slide templates with overflow/pagination.

        Args:
            template_path: Path to .pptx template file
            data: Data dict with oem_info, products, computed_fields, etc.
            output_path: Path for output .pptx file

        Returns:
            Result dict with generation details
        """
        start_time = time.time()
        errors = []
        shapes_populated = 0
        slides_created = 0

        try:
            prs = Presentation(template_path)

            # Calculate overflow slides needed
            overflow_slides_info = self._calculate_overflow_slides(data)

            # Create overflow slides if needed
            for overflow_info in overflow_slides_info:
                source_slide_idx = overflow_info['source_slide_index']
                num_overflow = overflow_info['overflow_count']
                for _ in range(num_overflow):
                    self._duplicate_slide(prs, source_slide_idx)
                    slides_created += 1

            # Get slide mappings grouped by slide_index
            slide_mappings = self._group_mappings_by_slide()

            # Process each slide
            total_slides = len(prs.slides)
            for slide_idx in range(total_slides):
                if slide_idx >= len(prs.slides):
                    break

                slide = prs.slides[slide_idx]

                # Get mappings for this slide (use slide 0 mappings for overflow slides)
                effective_slide_idx = 0 if slide_idx > 0 else slide_idx
                mappings = slide_mappings.get(effective_slide_idx, [])

                # Calculate product offset for this slide
                products = data.get("products", [])
                items_per_slide = self.pagination_rules.get("products", {}).get("items_per_slide", 2)
                product_offset = slide_idx * items_per_slide

                for shape_mapping in mappings:
                    try:
                        shape_id = shape_mapping["shape_id"]
                        shape = self._get_shape_by_id(slide, shape_id)

                        if shape is None:
                            if slide_idx == 0:  # Only report errors for first slide
                                errors.append(f"Shape {shape_id} not found in template")
                            continue

                        mapping_type = shape_mapping.get("type", "text")
                        is_repeatable = shape_mapping.get("repeatable_for") == "products"

                        if mapping_type == "text":
                            self._populate_text(shape, shape_mapping, data)
                            shapes_populated += 1

                        elif mapping_type == "table":
                            if shape.has_table:
                                if is_repeatable:
                                    # Products table - use offset for this slide
                                    self._populate_table(
                                        shape.table,
                                        shape_mapping,
                                        data,
                                        product_offset=product_offset
                                    )
                                else:
                                    self._populate_table(shape.table, shape_mapping, data)
                                shapes_populated += 1
                            else:
                                if slide_idx == 0:
                                    errors.append(f"Shape {shape_id} is not a table")

                        elif mapping_type == "bullet_list":
                            self._populate_bullet_list(shape, shape_mapping, data)
                            shapes_populated += 1

                    except Exception as e:
                        errors.append(f"Error populating shape {shape_mapping.get('shape_id')} on slide {slide_idx}: {str(e)}")

            # Save presentation
            Path(output_path).parent.mkdir(parents=True, exist_ok=True)
            prs.save(output_path)

            duration = time.time() - start_time

            return {
                "success": True,
                "presentation_path": output_path,
                "template_used": self.template_name,
                "shapes_populated": shapes_populated,
                "total_mappings": len(self.mapping.get("shape_mappings", [])),
                "slides_created": slides_created + 1,  # +1 for original slide
                "overflow_slides": slides_created,
                "errors": errors,
                "generation_timestamp": datetime.now().isoformat(),
                "generation_duration_seconds": round(duration, 2)
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "errors": errors,
                "generation_timestamp": datetime.now().isoformat()
            }

    def _calculate_overflow_slides(self, data: Dict) -> List[Dict]:
        """
        Calculate how many overflow slides are needed based on data.

        Returns:
            List of dicts with source_slide_index and overflow_count
        """
        overflow_info = []

        if not self.supports_multi_slide:
            return overflow_info

        for field_name, rules in self.pagination_rules.items():
            items = data.get(field_name, [])
            if not isinstance(items, list):
                continue

            items_per_slide = rules.get("items_per_slide", 2)
            source_slide_idx = rules.get("source_slide_index", 0)

            total_items = len(items)
            slides_needed = math.ceil(total_items / items_per_slide) if total_items > 0 else 1
            overflow_count = max(0, slides_needed - 1)

            if overflow_count > 0:
                overflow_info.append({
                    "field": field_name,
                    "source_slide_index": source_slide_idx,
                    "overflow_count": overflow_count,
                    "items_per_slide": items_per_slide
                })

        return overflow_info

    def _group_mappings_by_slide(self) -> Dict[int, List[Dict]]:
        """Group shape mappings by their slide_index."""
        grouped = {}
        for mapping in self.mapping.get("shape_mappings", []):
            slide_idx = mapping.get("slide_index", 0)
            if slide_idx not in grouped:
                grouped[slide_idx] = []
            grouped[slide_idx].append(mapping)
        return grouped

    def _duplicate_slide(self, prs: Presentation, source_idx: int) -> int:
        """
        Duplicate a slide and append it to the presentation.
        Uses python-pptx XML manipulation to copy the full slide layout.

        Args:
            prs: Presentation object
            source_idx: Index of slide to duplicate

        Returns:
            Index of the new slide
        """
        source_slide = prs.slides[source_idx]

        # Get the slide layout from source
        slide_layout = source_slide.slide_layout

        # Add a new slide with the same layout
        new_slide = prs.slides.add_slide(slide_layout)

        # Copy all shapes from source to new slide
        for shape in source_slide.shapes:
            self._copy_shape(shape, new_slide)

        return len(prs.slides) - 1

    def _copy_shape(self, shape, target_slide):
        """
        Copy a shape to a target slide.
        Handles different shape types appropriately.
        """
        # Get shape XML
        el = shape._element

        # Create a deep copy of the shape XML
        new_el = deepcopy(el)

        # Add to target slide's shape tree
        target_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    def _get_shape_by_id(self, slide, shape_id: int):
        """Find a shape by its shape_id."""
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        return None

    def _get_data_value(self, data: Dict, field_path: str, default: Any = None) -> Any:
        """
        Get a value from nested data using dot notation.

        Examples:
            "oem_info.oem_name" -> data["oem_info"]["oem_name"]
            "products[0].name" -> data["products"][0]["name"]
        """
        if not field_path:
            return default

        try:
            current = data
            parts = field_path.replace("][", ".").replace("[", ".").replace("]", "").split(".")

            for part in parts:
                if part.isdigit():
                    current = current[int(part)]
                else:
                    current = current.get(part, default) if isinstance(current, dict) else default
                    if current is None:
                        return default

            return current if current is not None else default
        except (KeyError, IndexError, TypeError):
            return default

    def _format_value(self, value: Any, format_str: str = "{value}", default: str = "N/A") -> str:
        """
        Format a value using format string.

        Args:
            value: The value to format
            format_str: Format string with {value} placeholder
            default: Default value if value is None/empty
        """
        if value is None or value == "" or value == 0:
            return default

        if isinstance(value, (int, float)):
            # Format numbers with thousand separators
            if value == int(value):
                formatted = f"{int(value):,}"
            else:
                formatted = f"{value:,.1f}"
        else:
            formatted = str(value)

        return format_str.replace("{value}", formatted)

    def _populate_text(self, shape, mapping: Dict, data: Dict):
        """Populate a text shape."""
        field = mapping.get("data_field", "")
        format_str = mapping.get("format", "{value}")
        default = mapping.get("default", "")

        value = self._get_data_value(data, field, default)
        formatted = self._format_value(value, format_str, default)

        if hasattr(shape, "text"):
            shape.text = formatted

    def _populate_table(self, table, mapping: Dict, data: Dict, product_offset: int = 0):
        """
        Populate a table shape.

        Args:
            table: Table object to populate
            mapping: Shape mapping configuration
            data: Data dictionary
            product_offset: Offset for products (used in multi-slide pagination)
        """
        row_mappings = mapping.get("row_mappings", [])
        product_columns = mapping.get("product_columns", [1, 2])
        max_products = mapping.get("max_products", 2)

        # Check if this is a products table (with array notation)
        is_products_table = any("products[]" in rm.get("data_field", "") for rm in row_mappings)

        if is_products_table:
            # Products table - populate multiple columns
            # Apply product_offset for multi-slide pagination
            all_products = data.get("products", [])
            products = all_products[product_offset:product_offset + max_products]

            for row_mapping in row_mappings:
                row_idx = row_mapping["row_index"]
                field_template = row_mapping.get("data_field", "")
                format_str = row_mapping.get("format", "{value}")
                default = row_mapping.get("default", "N/A")

                for prod_idx, col_idx in enumerate(product_columns):
                    if prod_idx >= len(products):
                        # Clear the cell if no product for this column
                        try:
                            cell = table.cell(row_idx, col_idx)
                            cell.text = ""
                        except IndexError:
                            pass
                        continue

                    # Replace products[] with products[idx] (using offset-adjusted index)
                    actual_idx = product_offset + prod_idx
                    field = field_template.replace("products[].", f"products.{actual_idx}.")
                    value = self._get_data_value(data, field, default)
                    formatted = self._format_value(value, format_str, default)

                    try:
                        cell = table.cell(row_idx, col_idx)
                        cell.text = formatted
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(8)
                    except IndexError:
                        pass  # Row/col out of bounds

        else:
            # Regular table - populate single column
            for row_mapping in row_mappings:
                row_idx = row_mapping["row_index"]
                col_idx = row_mapping.get("col_index", 1)
                field = row_mapping.get("data_field", "")
                format_str = row_mapping.get("format", "{value}")
                default = row_mapping.get("default", "")

                value = self._get_data_value(data, field, default)
                formatted = self._format_value(value, format_str, default)

                try:
                    cell = table.cell(row_idx, col_idx)
                    cell.text = formatted
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(9)
                except IndexError:
                    pass

    def _populate_bullet_list(self, shape, mapping: Dict, data: Dict):
        """Populate a shape with bullet points."""
        if not hasattr(shape, "text_frame"):
            return

        field = mapping.get("data_field", "")
        max_items = mapping.get("max_items", 5)
        font_size = mapping.get("font_size_pt", 9)

        items = self._get_data_value(data, field, [])

        if not isinstance(items, list):
            items = [items] if items else []

        items = items[:max_items]

        if not items:
            return

        tf = shape.text_frame

        # Clear existing paragraphs
        for p in tf.paragraphs:
            p.clear()

        # Add bullet items
        for i, item in enumerate(items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.level = 0
            run = p.add_run()
            run.text = str(item)
            run.font.size = Pt(font_size)


class DynamicOEMPresentationGenerator:
    """
    Higher-level generator that handles OEM data transformation and presentation generation.

    Similar interface to the original ppt_generator.py but uses dynamic mappings.
    """

    def __init__(self, mapping_path: str = "src/config/template_schemas/iaa_template.json"):
        """Initialize with a mapping configuration."""
        self.mapping_path = mapping_path
        with open(mapping_path, 'r') as f:
            self.mapping = json.load(f)
        self.generator = DynamicPPTGenerator(self.mapping)

    def generate_from_scraping_result(
        self,
        scraping_result: Dict,
        template_path: str = "templates/IAA_Template.pptx",
        output_dir: str = "outputs",
        max_products: Optional[int] = None
    ) -> Dict[str, Any]:
        """
        Generate presentation from scraping result.
        Supports multi-slide output when there are more products than fit on one slide.

        Args:
            scraping_result: ScrapingResult dict from scraper
            template_path: Path to template file
            output_dir: Output directory
            max_products: Maximum products to include (None = all)

        Returns:
            Result dict with presentation details
        """
        # Transform scraping result to template data format
        data = self._transform_scraping_result(scraping_result, max_products)

        # Generate safe filename
        oem_name = data.get("oem_info", {}).get("oem_name", "Unknown")
        safe_name = re.sub(r'[^\w\-]', '_', oem_name)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"{output_dir}/Dynamic_{safe_name}_{timestamp}.pptx"

        # Generate presentation
        result = self.generator.generate(template_path, data, output_path)

        return {
            "oem_name": oem_name,
            "products_count": len(data.get("products", [])),
            "data": data,
            "presentation_result": result
        }

    def _transform_scraping_result(self, scraping_result: Dict, max_products: Optional[int] = None) -> Dict:
        """
        Transform scraping result to template data format.

        Args:
            scraping_result: Raw scraping result
            max_products: Maximum products to include (None = all)

        Returns:
            Transformed data dict for template generation
        """
        oem_name = scraping_result.get('oem_name', 'Unknown OEM')
        oem_url = scraping_result.get('oem_url', '')
        vehicles = scraping_result.get('vehicles', [])

        # Extract domain
        domain = oem_url.split('/')[2] if '://' in oem_url else oem_url

        # Transform vehicles to products (no limit by default for multi-slide support)
        vehicles_to_transform = vehicles if max_products is None else vehicles[:max_products]
        products = [self._transform_vehicle(v) for v in vehicles_to_transform]

        # Generate computed fields
        highlights = self._generate_highlights(vehicles, oem_name)
        assessment = self._generate_assessment(scraping_result)
        technologies = self._extract_technologies(vehicles)

        return {
            "oem_info": {
                "oem_name": oem_name,
                "country": self._extract_country(oem_url),
                "address": "",
                "website": domain,
                "booth": "",
                "category": self._determine_category(vehicles)
            },
            "products": products,
            "computed_fields": {
                "expected_highlights": highlights,
                "assessment": assessment,
                "technologies": technologies
            },
            "source_url": oem_url,
            "extraction_timestamp": scraping_result.get('extraction_timestamp', datetime.now().isoformat())
        }

    def _transform_vehicle(self, vehicle: Dict) -> Dict:
        """Transform a vehicle to product format."""
        def fmt(value, unit=""):
            if value is None or value == 0:
                return "N/A"
            if isinstance(value, (int, float)):
                return f"{value:,.0f} {unit}".strip() if value == int(value) else f"{value:,.1f} {unit}".strip()
            return str(value) if value else "N/A"

        additional = vehicle.get('additional_specs', {}) or {}

        return {
            "name": vehicle.get('vehicle_name', 'Unknown Vehicle'),
            "wheel_formula": additional.get('wheel_formula', 'N/A'),
            "wheelbase": fmt(additional.get('wheelbase_mm'), 'mm'),
            "gvw_gcw": fmt(vehicle.get('gvw_kg'), 'kg GVW'),
            "range": fmt(vehicle.get('range_km'), 'km'),
            "battery": fmt(vehicle.get('battery_capacity_kwh'), 'kWh'),
            "fuel_cell": fmt(additional.get('fuel_cell_kw'), 'kW'),
            "h2_tank": fmt(additional.get('h2_tank_kg'), 'kg'),
            "charging": fmt(vehicle.get('dc_charging_kw'), 'kW DC'),
            "performance": fmt(vehicle.get('motor_power_kw'), 'kW'),
            "powertrain": additional.get('powertrain_description', vehicle.get('powertrain_type', 'Electric')),
            "sop": additional.get('start_of_production', 'TBD'),
            "markets": additional.get('markets', 'N/A'),
            "application": vehicle.get('category', 'Commercial Vehicle')
        }

    def _generate_highlights(self, vehicles: List[Dict], oem_name: str) -> List[str]:
        """Generate highlights from vehicle data."""
        highlights = []

        if vehicles:
            bev_count = sum(1 for v in vehicles if v.get('battery_capacity_kwh'))
            if bev_count > 0:
                highlights.append(f"{bev_count} Battery Electric Vehicle(s) in portfolio")

            max_range = max((v.get('range_km', 0) or 0 for v in vehicles), default=0)
            if max_range > 0:
                highlights.append(f"Up to {max_range:.0f} km range capability")

            max_battery = max((v.get('battery_capacity_kwh', 0) or 0 for v in vehicles), default=0)
            if max_battery > 0:
                highlights.append(f"Battery capacity up to {max_battery:.0f} kWh")

        if not highlights:
            highlights = [f"{oem_name} e-mobility portfolio", "Technical specifications available"]

        return highlights[:4]

    def _generate_assessment(self, scraping_result: Dict) -> List[str]:
        """Generate assessment from scraping result."""
        assessment = []
        score = scraping_result.get('source_compliance_score', 0)
        vehicles = scraping_result.get('vehicles', [])

        if score >= 0.8:
            assessment.append("High data quality from official sources")
        elif score >= 0.5:
            assessment.append("Moderate data quality, some gaps")
        else:
            assessment.append("Limited data available")

        if len(vehicles) >= 1:
            assessment.append(f"{len(vehicles)} vehicle(s) documented")

        return assessment[:3]

    def _extract_technologies(self, vehicles: List[Dict]) -> List[str]:
        """Extract technologies from vehicles."""
        technologies = set()

        for v in vehicles:
            if v.get('battery_capacity_kwh'):
                technologies.add("Li-ion Battery")

            dc_power = v.get('dc_charging_kw')
            if dc_power and dc_power >= 150:
                technologies.add(f"High-Power DC Charging ({dc_power:.0f}kW)")
            elif dc_power:
                technologies.add("DC Fast Charging")

        return list(technologies)[:5]

    def _extract_country(self, url: str) -> str:
        """Extract country from URL."""
        tld_map = {'.de': 'Germany', '.eu': 'Europe', '.com': 'USA', '.cn': 'China'}
        for tld, country in tld_map.items():
            if tld in url:
                return country
        return ""

    def _determine_category(self, vehicles: List[Dict]) -> str:
        """Determine OEM category."""
        categories = [v.get('category', '').lower() for v in vehicles]
        if any('truck' in c for c in categories):
            return "OEM - Commercial Trucks"
        elif any('bus' in c for c in categories):
            return "OEM - Bus & Coach"
        return "OEM - Commercial Vehicles"


# =============================================================================
# CLI for testing
# =============================================================================

if __name__ == "__main__":
    import sys

    print("=" * 60)
    print("Dynamic PPT Generator Test - Multi-Slide Support")
    print("=" * 60)

    # Test with mock data - 5 products to test multi-slide overflow
    mock_scraping_result = {
        "oem_name": "Test OEM (Multi-Slide)",
        "oem_url": "https://www.testoem.eu",
        "vehicles": [
            {
                "vehicle_name": "eTruck X1",
                "category": "Heavy-duty Truck",
                "powertrain_type": "BEV",
                "battery_capacity_kwh": 500,
                "range_km": 400,
                "motor_power_kw": 350,
                "dc_charging_kw": 350,
                "gvw_kg": 40000,
                "additional_specs": {
                    "wheel_formula": "6x2",
                    "wheelbase_mm": 4500,
                    "start_of_production": "2025",
                    "markets": "EU"
                }
            },
            {
                "vehicle_name": "eTruck X2",
                "category": "Heavy-duty Truck",
                "powertrain_type": "BEV",
                "battery_capacity_kwh": 600,
                "range_km": 500,
                "motor_power_kw": 400,
                "dc_charging_kw": 400,
                "gvw_kg": 44000,
                "additional_specs": {
                    "wheel_formula": "6x4",
                    "wheelbase_mm": 4800,
                    "start_of_production": "2025",
                    "markets": "EU, NA"
                }
            },
            {
                "vehicle_name": "eTruck X3 Long Range",
                "category": "Heavy-duty Truck",
                "powertrain_type": "BEV",
                "battery_capacity_kwh": 800,
                "range_km": 700,
                "motor_power_kw": 450,
                "dc_charging_kw": 450,
                "gvw_kg": 40000,
                "additional_specs": {
                    "wheel_formula": "6x2",
                    "wheelbase_mm": 5000,
                    "start_of_production": "2026",
                    "markets": "Global"
                }
            },
            {
                "vehicle_name": "eBus City 12m",
                "category": "City Bus",
                "powertrain_type": "BEV",
                "battery_capacity_kwh": 350,
                "range_km": 300,
                "motor_power_kw": 250,
                "dc_charging_kw": 150,
                "gvw_kg": 19000,
                "additional_specs": {
                    "wheel_formula": "4x2",
                    "wheelbase_mm": 5900,
                    "start_of_production": "2024",
                    "markets": "EU"
                }
            },
            {
                "vehicle_name": "eBus Coach 13m",
                "category": "Coach",
                "powertrain_type": "BEV",
                "battery_capacity_kwh": 450,
                "range_km": 400,
                "motor_power_kw": 300,
                "dc_charging_kw": 250,
                "gvw_kg": 24000,
                "additional_specs": {
                    "wheel_formula": "6x2",
                    "wheelbase_mm": 6500,
                    "start_of_production": "2025",
                    "markets": "EU, APAC"
                }
            }
        ],
        "source_compliance_score": 0.85,
        "extraction_timestamp": datetime.now().isoformat()
    }

    try:
        # Test the generator with all 5 products (should create 3 slides: 2+2+1 products)
        print("\n--- Test 1: Multi-Slide (5 products) ---")
        generator = DynamicOEMPresentationGenerator()
        result = generator.generate_from_scraping_result(mock_scraping_result)

        print(f"\nGeneration Result:")
        print(f"  OEM: {result['oem_name']}")
        print(f"  Products: {result['products_count']}")
        print(f"  Success: {result['presentation_result'].get('success')}")
        print(f"  Output: {result['presentation_result'].get('presentation_path')}")
        print(f"  Slides created: {result['presentation_result'].get('slides_created', 1)}")
        print(f"  Overflow slides: {result['presentation_result'].get('overflow_slides', 0)}")
        print(f"  Shapes populated: {result['presentation_result'].get('shapes_populated')}")
        print(f"  Duration: {result['presentation_result'].get('generation_duration_seconds')}s")

        if result['presentation_result'].get('errors'):
            print(f"  Errors: {result['presentation_result']['errors']}")

        # Test 2: Single slide (2 products)
        print("\n--- Test 2: Single Slide (2 products, limited) ---")
        result2 = generator.generate_from_scraping_result(mock_scraping_result, max_products=2)
        print(f"  Products: {result2['products_count']}")
        print(f"  Slides created: {result2['presentation_result'].get('slides_created', 1)}")
        print(f"  Success: {result2['presentation_result'].get('success')}")

    except FileNotFoundError as e:
        print(f"ERROR: {e}")
        print("Make sure template and mapping files exist.")
    except Exception as e:
        print(f"ERROR: {e}")
        import traceback
        traceback.print_exc()
