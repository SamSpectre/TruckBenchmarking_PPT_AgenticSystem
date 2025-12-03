"""
Template Analyzer for Plug-and-Play PowerPoint System

Discovers the structure of any PowerPoint template by analyzing:
- Shapes (text boxes, placeholders, images)
- Tables (dimensions, headers, row labels)
- Positions and dimensions

This is a standalone module - does NOT modify existing ppt_generator.py
"""

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.base import BaseShape
from pptx.shapes.placeholder import PlaceholderPicture
from typing import Dict, List, Any, Optional
from pathlib import Path
from datetime import datetime
import hashlib
import json


class TemplateAnalyzer:
    """
    Analyzes PowerPoint templates to discover their structure.

    Usage:
        analyzer = TemplateAnalyzer()
        structure = analyzer.analyze("path/to/template.pptx")
        print(json.dumps(structure, indent=2))
    """

    def __init__(self):
        self.emu_to_inches = lambda emu: round(emu / 914400, 2) if emu else 0

    def analyze(self, template_path: str) -> Dict[str, Any]:
        """
        Analyze a PowerPoint template and return its structure.

        Args:
            template_path: Path to the .pptx file

        Returns:
            Dict containing template structure with slides, shapes, tables
        """
        path = Path(template_path)
        if not path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")

        if not path.suffix.lower() == '.pptx':
            raise ValueError(f"File must be .pptx format: {template_path}")

        # Load presentation
        prs = Presentation(template_path)

        # Calculate template hash for caching/versioning
        with open(template_path, 'rb') as f:
            template_hash = hashlib.md5(f.read()).hexdigest()[:12]

        # Build structure
        structure = {
            "template_name": path.stem,
            "template_path": str(path.absolute()),
            "template_hash": template_hash,
            "analyzed_at": datetime.now().isoformat(),
            "slide_width_inches": self.emu_to_inches(prs.slide_width),
            "slide_height_inches": self.emu_to_inches(prs.slide_height),
            "total_slides": len(prs.slides),
            "slides": []
        }

        # Analyze each slide
        for slide_idx, slide in enumerate(prs.slides):
            slide_info = self._analyze_slide(slide, slide_idx)
            structure["slides"].append(slide_info)

        # Generate summary
        structure["summary"] = self._generate_summary(structure)

        return structure

    def _analyze_slide(self, slide, slide_index: int) -> Dict[str, Any]:
        """Analyze a single slide and its shapes."""
        slide_info = {
            "slide_index": slide_index,
            "slide_id": slide.slide_id if hasattr(slide, 'slide_id') else None,
            "shapes": [],
            "tables": [],
            "placeholders": [],
            "text_boxes": [],
            "images": []
        }

        for shape in slide.shapes:
            shape_info = self._analyze_shape(shape)

            # Categorize by type
            slide_info["shapes"].append(shape_info)

            if shape_info["type"] == "TABLE":
                slide_info["tables"].append(shape_info)
            elif shape_info["type"] == "PLACEHOLDER":
                slide_info["placeholders"].append(shape_info)
            elif shape_info["type"] == "TEXT_BOX":
                slide_info["text_boxes"].append(shape_info)
            elif shape_info["type"] in ["PICTURE", "PLACEHOLDER_PICTURE"]:
                slide_info["images"].append(shape_info)

        return slide_info

    def _analyze_shape(self, shape: BaseShape) -> Dict[str, Any]:
        """Analyze a single shape and extract relevant info."""
        shape_info = {
            "shape_id": shape.shape_id,
            "name": shape.name,
            "type": self._get_shape_type(shape),
            "position": {
                "left_inches": self.emu_to_inches(shape.left),
                "top_inches": self.emu_to_inches(shape.top),
                "width_inches": self.emu_to_inches(shape.width),
                "height_inches": self.emu_to_inches(shape.height)
            }
        }

        # Extract text content if available
        if hasattr(shape, "text_frame"):
            shape_info["has_text_frame"] = True
            shape_info["text_content"] = self._extract_text(shape)
            shape_info["paragraph_count"] = len(shape.text_frame.paragraphs)
        else:
            shape_info["has_text_frame"] = False

        # Handle placeholders
        if shape.is_placeholder:
            shape_info["is_placeholder"] = True
            shape_info["placeholder_type"] = self._get_placeholder_type(shape)
            shape_info["placeholder_idx"] = shape.placeholder_format.idx if hasattr(shape, 'placeholder_format') else None
        else:
            shape_info["is_placeholder"] = False

        # Handle tables
        if shape.has_table:
            shape_info["table_info"] = self._analyze_table(shape.table)

        return shape_info

    def _get_shape_type(self, shape: BaseShape) -> str:
        """Determine the type of shape."""
        if shape.has_table:
            return "TABLE"

        if shape.is_placeholder:
            if isinstance(shape, PlaceholderPicture):
                return "PLACEHOLDER_PICTURE"
            return "PLACEHOLDER"

        # Check shape type enum
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                return "TEXT_BOX"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return "PICTURE"
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return "AUTO_SHAPE"
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                return "GROUP"
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                return "CHART"
            else:
                return str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
        except:
            return "UNKNOWN"

    def _get_placeholder_type(self, shape: BaseShape) -> Optional[str]:
        """Get the placeholder type name."""
        try:
            ph_type = shape.placeholder_format.type
            return str(ph_type).replace("PP_PLACEHOLDER.", "")
        except:
            return None

    def _extract_text(self, shape: BaseShape) -> str:
        """Extract all text from a shape."""
        try:
            if hasattr(shape, "text"):
                return shape.text.strip()
            return ""
        except:
            return ""

    def _analyze_table(self, table) -> Dict[str, Any]:
        """Analyze a table's structure."""
        table_info = {
            "rows": len(table.rows),
            "columns": len(table.columns),
            "cells": []
        }

        # Extract cell contents
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                cell_text = cell.text.strip() if cell.text else ""
                if cell_text:  # Only include non-empty cells
                    table_info["cells"].append({
                        "row": row_idx,
                        "col": col_idx,
                        "text": cell_text
                    })

        # Try to identify headers (first row)
        if len(table.rows) > 0:
            headers = []
            for cell in table.rows[0].cells:
                headers.append(cell.text.strip() if cell.text else "")
            table_info["headers"] = headers

        # Try to identify row labels (first column)
        if len(table.columns) > 0:
            row_labels = []
            for row in table.rows:
                if len(row.cells) > 0:
                    row_labels.append(row.cells[0].text.strip() if row.cells[0].text else "")
            table_info["row_labels"] = row_labels

        return table_info

    def _generate_summary(self, structure: Dict) -> Dict[str, Any]:
        """Generate a summary of the template structure."""
        total_shapes = 0
        total_tables = 0
        total_placeholders = 0
        total_text_boxes = 0

        for slide in structure["slides"]:
            total_shapes += len(slide["shapes"])
            total_tables += len(slide["tables"])
            total_placeholders += len(slide["placeholders"])
            total_text_boxes += len(slide["text_boxes"])

        return {
            "total_shapes": total_shapes,
            "total_tables": total_tables,
            "total_placeholders": total_placeholders,
            "total_text_boxes": total_text_boxes,
            "has_tables": total_tables > 0,
            "has_placeholders": total_placeholders > 0
        }

    def analyze_for_mapping(self, template_path: str) -> Dict[str, Any]:
        """
        Analyze template and return a simplified structure optimized for LLM mapping.

        This returns a cleaner format that's easier for the LLM to understand
        and map to vehicle data fields.
        """
        full_analysis = self.analyze(template_path)

        # Simplified structure for LLM
        mapping_ready = {
            "template_name": full_analysis["template_name"],
            "template_hash": full_analysis["template_hash"],
            "mappable_elements": []
        }

        for slide in full_analysis["slides"]:
            for shape in slide["shapes"]:
                element = {
                    "slide_index": slide["slide_index"],
                    "shape_id": shape["shape_id"],
                    "shape_name": shape["name"],
                    "type": shape["type"],
                    "current_content": shape.get("text_content", "")
                }

                # Add table structure if applicable
                if "table_info" in shape:
                    element["table_structure"] = {
                        "rows": shape["table_info"]["rows"],
                        "columns": shape["table_info"]["columns"],
                        "headers": shape["table_info"].get("headers", []),
                        "row_labels": shape["table_info"].get("row_labels", [])
                    }

                # Add placeholder info
                if shape.get("is_placeholder"):
                    element["placeholder_type"] = shape.get("placeholder_type")

                mapping_ready["mappable_elements"].append(element)

        return mapping_ready

    def save_analysis(self, structure: Dict, output_path: str):
        """Save analysis to JSON file."""
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(structure, f, indent=2, ensure_ascii=False)
        print(f"Analysis saved to: {output_path}")


# =============================================================================
# CLI for standalone testing
# =============================================================================

if __name__ == "__main__":
    import sys

    print("=" * 60)
    print("PowerPoint Template Analyzer")
    print("=" * 60)

    # Default to IAA template if no argument
    if len(sys.argv) > 1:
        template_path = sys.argv[1]
    else:
        template_path = "templates/IAA_Template.pptx"

    print(f"\nAnalyzing: {template_path}")

    analyzer = TemplateAnalyzer()

    try:
        # Full analysis
        structure = analyzer.analyze(template_path)

        print(f"\n--- Template Summary ---")
        print(f"Name: {structure['template_name']}")
        print(f"Hash: {structure['template_hash']}")
        print(f"Slides: {structure['total_slides']}")
        print(f"Total shapes: {structure['summary']['total_shapes']}")
        print(f"Tables: {structure['summary']['total_tables']}")
        print(f"Placeholders: {structure['summary']['total_placeholders']}")
        print(f"Text boxes: {structure['summary']['total_text_boxes']}")

        # Show shapes on first slide
        if structure["slides"]:
            print(f"\n--- Slide 0 Shapes ---")
            for shape in structure["slides"][0]["shapes"]:
                text_preview = shape.get("text_content", "")[:30] + "..." if shape.get("text_content", "") else ""
                print(f"  ID {shape['shape_id']:2d} | {shape['type']:15s} | {shape['name']:20s} | {text_preview}")

        # Save full analysis
        output_path = f"outputs/{structure['template_name']}_analysis.json"
        analyzer.save_analysis(structure, output_path)

        # Also show mapping-ready format
        mapping_ready = analyzer.analyze_for_mapping(template_path)
        print(f"\n--- Mapping-Ready Elements ({len(mapping_ready['mappable_elements'])}) ---")
        for elem in mapping_ready["mappable_elements"][:5]:
            print(f"  Shape {elem['shape_id']}: {elem['type']} - '{elem['current_content'][:40]}'")

    except FileNotFoundError as e:
        print(f"ERROR: {e}")
    except Exception as e:
        print(f"ERROR: {e}")
        import traceback
        traceback.print_exc()
